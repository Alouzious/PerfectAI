"""
File Processor Service
Extracts slides from PDF and PPTX files
"""
import os
from pptx import Presentation
from PyPDF2 import PdfReader
from pdf2image import convert_from_path
import logging

logger = logging.getLogger(__name__)


class FileProcessor:
    """Process pitch deck files and extract slide content"""
    
    def __init__(self):
        self.supported_formats = ['pdf', 'pptx', 'ppt']
    
    def extract_slides(self, file_path):
        """
        Main method to extract slides from any supported file type
        
        Args:
            file_path (str): Path to the uploaded file
            
        Returns:
            list: List of slide dictionaries with extracted content
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        file_extension = file_path.split('.')[-1].lower()
        
        if file_extension in ['pptx', 'ppt']:
            return self._extract_from_pptx(file_path)
        elif file_extension == 'pdf':
            return self._extract_from_pdf(file_path)
        else:
            raise ValueError(f"Unsupported file format: {file_extension}")
    
    def _extract_from_pptx(self, file_path):
        """
        Extract content from PowerPoint files
        
        Args:
            file_path (str): Path to PPTX file
            
        Returns:
            list: Slide data
        """
        slides_data = []
        
        try:
            prs = Presentation(file_path)
            logger.info(f"Processing PPTX with {len(prs.slides)} slides")
            
            for idx, slide in enumerate(prs.slides, start=1):
                text_content = []
                has_images = False
                has_charts = False
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    # Text content
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text.strip())
                    
                    # Check for images
                    if shape.shape_type == 13:  # Picture
                        has_images = True
                    
                    # Check for charts
                    if hasattr(shape, "has_chart") and shape.has_chart:
                        has_charts = True
                
                # Extract notes (speaker notes)
                notes_text = ""
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if notes_slide.notes_text_frame:
                        notes_text = notes_slide.notes_text_frame.text.strip()
                
                slide_data = {
                    'number': idx,
                    'text': '\n'.join(text_content),
                    'notes': notes_text,
                    'has_images': has_images,
                    'has_charts': has_charts,
                    'word_count': len(' '.join(text_content).split()),
                }
                
                slides_data.append(slide_data)
                logger.info(f"Extracted slide {idx}: {len(text_content)} text blocks")
            
            return slides_data
            
        except Exception as e:
            logger.error(f"Error extracting from PPTX: {str(e)}")
            raise
    
    def _extract_from_pdf(self, file_path):
        """
        Extract content from PDF files
        
        Args:
            file_path (str): Path to PDF file
            
        Returns:
            list: Slide data
        """
        slides_data = []
        
        try:
            reader = PdfReader(file_path)
            total_pages = len(reader.pages)
            logger.info(f"Processing PDF with {total_pages} pages")
            
            for idx, page in enumerate(reader.pages, start=1):
                # Extract text
                text = page.extract_text()
                
                # Check for images (basic check)
                has_images = False
                if '/XObject' in page['/Resources']:
                    xObject = page['/Resources']['/XObject'].get_object()
                    for obj in xObject:
                        if xObject[obj]['/Subtype'] == '/Image':
                            has_images = True
                            break
                
                slide_data = {
                    'number': idx,
                    'text': text.strip() if text else '',
                    'notes': '',  # PDFs don't have speaker notes
                    'has_images': has_images,
                    'has_charts': False,  # Hard to detect in PDFs
                    'word_count': len(text.split()) if text else 0,
                }
                
                slides_data.append(slide_data)
                logger.info(f"Extracted page {idx}")
            
            return slides_data
            
        except Exception as e:
            logger.error(f"Error extracting from PDF: {str(e)}")
            raise
    
    def get_file_metadata(self, file_path):
        """
        Get metadata about the file
        
        Args:
            file_path (str): Path to file
            
        Returns:
            dict: File metadata
        """
        file_size = os.path.getsize(file_path)
        file_extension = file_path.split('.')[-1].lower()
        
        return {
            'size_bytes': file_size,
            'size_mb': round(file_size / (1024 * 1024), 2),
            'extension': file_extension,
        }